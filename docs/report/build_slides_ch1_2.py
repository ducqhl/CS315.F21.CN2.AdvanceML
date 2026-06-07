#!/usr/bin/env python3
"""
Build Chapters 1-2 slides for the Crypto Lambda-Architecture project,
reusing the visual template of BaoCaoDoAn.AdvancedML.pptx (House-Price deck).

Strategy: copy the template, rewrite slide text IN-PLACE (preserving the exact
fonts/colors/bullets/layouts), trim to chapters 1-2, then duplicate a clean
bullet-list slide for the related-work topics and drop in two repo figures.
"""
import copy
import os
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, "BaoCaoDoAn.AdvancedML.pptx")
OUT = os.path.join(HERE, "BaoCaoDoAn.AdvancedML.Ch1-2.pptx")
FIG = os.path.join(HERE, "figures")

prs = Presentation(SRC)


# ---------- helpers ----------------------------------------------------------
def get_shape(slide, shape_id):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    raise KeyError(f"shape {shape_id} not found")


def _make_run(rpr_template, text, bold):
    r = etree.SubElement(etree.Element("dummy"), qn("a:r"))
    # build via a fresh element to control child order: rPr then t
    r = parse_xml('<a:r xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
    if rpr_template is not None:
        rpr = copy.deepcopy(rpr_template)
    else:
        rpr = parse_xml('<a:rPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
    if bold is not None:
        rpr.set("b", "1" if bold else "0")
    r.append(rpr)
    t = parse_xml('<a:t xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
    t.text = text
    r.append(t)
    return r


def set_lines(shape, lines):
    """Replace the shape's paragraphs with `lines`.

    Each line is either a str, or a list of (text, bold|None) segments.
    Paragraph properties (bullets/indent/align) come from the template's
    first paragraph; each segment reuses the rPr of the correspondingly
    positioned original run (so a bold-label + light-value pattern survives),
    falling back to the first run's rPr.
    """
    tf = shape.text_frame
    txBody = tf._txBody
    p_els = txBody.findall(qn("a:p"))
    template_p = copy.deepcopy(p_els[0])
    template_ppr = template_p.find(qn("a:pPr"))
    template_rprs = [copy.deepcopy(r.find(qn("a:rPr")))
                     for r in template_p.findall(qn("a:r"))
                     if r.find(qn("a:rPr")) is not None]
    template_endpr = template_p.find(qn("a:endParaRPr"))

    for p in p_els:
        txBody.remove(p)

    for line in lines:
        new_p = parse_xml('<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
        if template_ppr is not None:
            new_p.append(copy.deepcopy(template_ppr))
        segments = [(line, None)] if isinstance(line, str) else line
        for i, (text, bold) in enumerate(segments):
            base = template_rprs[i] if i < len(template_rprs) else (
                template_rprs[0] if template_rprs else None)
            new_p.append(_make_run(base, text, bold))
        if template_endpr is not None:
            new_p.append(copy.deepcopy(template_endpr))
        txBody.append(new_p)


def duplicate_slide(prs, index):
    source = prs.slides[index]
    layout = source.slide_layout
    new_slide = prs.slides.add_slide(layout)
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in source.shapes:
        new_slide.shapes._spTree.append(copy.deepcopy(shp._element))
    return new_slide


# ============================================================================
# CHAPTER 1 : GIỚI THIỆU ĐỀ TÀI  (slides 0-8, edited in place)
# ============================================================================

# --- Slide 0 : Title ---
s = prs.slides[0]
set_lines(get_shape(s, 1232), ["Báo cáo đồ án môn học", "Máy học nâng cao (CS315)"])
set_lines(get_shape(s, 1233),
          ["Hệ thống phân tích và dự đoán giá tiền mã hóa theo thời gian thực dựa trên kiến trúc Lambda"])
set_lines(get_shape(s, 1234), [
    [("GVHD: ", True), ("[Điền tên giảng viên]", False)],
    [("Lớp: ", True), ("CS315.F21.CN2", False)],
    [("Nhóm thực hiện: ", True), ("Nhóm CN2", False)],
])

# --- Slide 1 : Team ---
s = prs.slides[1]
set_lines(get_shape(s, 1239), ["Nhân sự nhóm CN2"])
set_lines(get_shape(s, 1240), [
    "Thành viên 1 — MSSV",
    "Thành viên 2 — MSSV",
    "Thành viên 3 — MSSV",
    "Thành viên 4 — MSSV",
])

# --- Slide 2 : Agenda (6 report chapters) ---
s = prs.slides[2]
set_lines(get_shape(s, 1245), ["Nội dung thuyết trình"])
set_lines(get_shape(s, 1255), ["Giới thiệu đề tài"])
set_lines(get_shape(s, 1261), ["Bối cảnh, bài toán & mục tiêu nghiên cứu"])
set_lines(get_shape(s, 1256), ["Nghiên cứu liên quan"])
set_lines(get_shape(s, 1262), ["Các hướng tiếp cận & điểm mới của đề tài"])
set_lines(get_shape(s, 1257), ["Dữ liệu & phân tích khám phá"])
set_lines(get_shape(s, 1263), ["Thu thập, EDA và xây dựng đặc trưng"])
set_lines(get_shape(s, 1258), ["Phương pháp đề xuất"])
set_lines(get_shape(s, 1246), ["Kiến trúc Lambda & mô hình LSTM v3"])
set_lines(get_shape(s, 1259), ["Thực nghiệm & kết quả"])
set_lines(get_shape(s, 1247), ["Đánh giá, walk-forward & backtest"])
set_lines(get_shape(s, 1260), ["Kết luận"])
set_lines(get_shape(s, 1248), ["Đóng góp, hạn chế & hướng phát triển"])

# --- Slide 3 : Section header 01 ---
s = prs.slides[3]
set_lines(get_shape(s, 1268), ["Giới thiệu đề tài"])
set_lines(get_shape(s, 1269), ["01"])

# --- Slide 4 : Tính cấp thiết của đề tài ---
s = prs.slides[4]
set_lines(get_shape(s, 1275), ["Tính cấp thiết của đề tài"])
set_lines(get_shape(s, 1276), [
    [("Thị trường crypto ", True),
     ("biến động cao — dự đoán xu hướng ngắn hạn mang giá trị lớn cho nhà đầu tư", False)],
    [("Hạn chế hiện tại: ", True),
     ("phần lớn hệ thống chỉ xử lý batch, thiếu phân tích gần thời gian thực", False)],
])

# --- Slide 5 : Mô tả bài toán ---
s = prs.slides[5]
set_lines(get_shape(s, 1281), ["Mô tả bài toán"])
set_lines(get_shape(s, 1282), [
    [("Input: ", True), ("giá thời gian thực BTC & DOGE từ CoinGecko (giá, khối lượng, vốn hóa)", False)],
    [("Output 1: ", True), ("chỉ số kỹ thuật thời gian thực — SMA, RSI, Bollinger, VWAP, ATR", False)],
    [("Output 2: ", True), ("dự đoán giá 7 ngày kèm hướng xu hướng (tăng/giảm)", False)],
    [("Output 3: ", True), ("dashboard trực quan hóa toàn bộ pipeline", False)],
])

# --- Slide 6 : Đối tượng & phạm vi ---
s = prs.slides[6]
set_lines(get_shape(s, 1287), ["Đối tượng & phạm vi nghiên cứu"])
set_lines(get_shape(s, 1288), [
    [("Đối tượng: ", True), ("Bitcoin (store-of-value) & Dogecoin (tính đầu cơ/meme cao)", False)],
    [("Dữ liệu: ", True), ("lịch sử CoinGecko ~11,4 năm (2015–2026), tối thiểu 2 năm", False)],
    [("Dự đoán: ", True), ("giá ngắn hạn 7 ngày tiếp theo", False)],
    [("Triển khai: ", True), ("cục bộ bằng Docker Compose (9 services)", False)],
])

# --- Slide 7 : Mục tiêu đề tài ---
s = prs.slides[7]
set_lines(get_shape(s, 1293), ["Mục tiêu đề tài"])
set_lines(get_shape(s, 1294), [
    "Xây dựng pipeline dữ liệu Lambda Architecture hoàn chỉnh",
    "Tính toán chỉ số kỹ thuật theo thời gian thực bằng Spark Streaming",
    "Huấn luyện mô hình LSTM dual-head dự đoán giá và xu hướng",
    "Xây dựng API (FastAPI) và giao diện người dùng (React + Streamlit)",
])

# --- Slide 8 : Kiến trúc tổng quát (Lambda) — 4-step process ---
s = prs.slides[8]
set_lines(get_shape(s, 1299), ["Kiến trúc tổng quát hệ thống (Lambda)"])
set_lines(get_shape(s, 1304), ["1. Thu thập dữ liệu"])
set_lines(get_shape(s, 1300), ["CoinGecko API → Kafka (topic crypto_raw), polling mỗi 10 phút"])
set_lines(get_shape(s, 1306), ["2. Xử lý Batch & Streaming"])
set_lines(get_shape(s, 1301), ["Spark Streaming tính chỉ số kỹ thuật; Spark Batch tổng hợp lịch sử"])
set_lines(get_shape(s, 1305), ["3. Lưu trữ & phục vụ"])
set_lines(get_shape(s, 1302), ["MongoDB (serving layer) → FastAPI REST API (JWT auth)"])
set_lines(get_shape(s, 1307), ["4. Dự đoán & trực quan"])
set_lines(get_shape(s, 1303), ["LSTM dự đoán 7 ngày → React + Streamlit dashboard"])

# ============================================================================
# CHAPTER 2 : NGHIÊN CỨU LIÊN QUAN  (slides 9-10 in place, then 5 duplicates)
# ============================================================================

# --- Slide 9 : Section header 02 ---
s = prs.slides[9]
set_lines(get_shape(s, 1312), ["Nghiên cứu liên quan"])
set_lines(get_shape(s, 1313), ["02"])

# --- Slide 10 : Tổng quan các hướng nghiên cứu (CUSTOM_6, 3 blocks) ---
s = prs.slides[10]
set_lines(get_shape(s, 1319), ["Tổng quan các hướng nghiên cứu"])
set_lines(get_shape(s, 1323), ["Học sâu (LSTM)"])
set_lines(get_shape(s, 1320), ["Nắm bắt phụ thuộc dài hạn; vượt mô hình thống kê truyền thống [Fischer & Krauss, 2018]"])
set_lines(get_shape(s, 1324), ["Xử lý thời gian thực"])
set_lines(get_shape(s, 1321), [
    "Kafka + Spark Structured Streaming [Zaharia et al., 2016]",
    "Cửa sổ trượt, watermark, ghi MongoDB qua foreachBatch",
])
set_lines(get_shape(s, 1325), ["Kiến trúc Lambda"])
set_lines(get_shape(s, 1322), [
    "Batch layer: chính xác cao, độ trễ cao",
    "Speed layer: độ trễ thấp, kết quả xấp xỉ",
    "Serving layer: phục vụ truy vấn [Marz & Warren, 2015]",
    "Phù hợp hệ thống tài chính: lịch sử + realtime",
])

# --- Trim deck to chapters 1-2: remove everything after slide 10 ---
# Drop both the sldId entry AND the presentation->slide relationship so the
# orphaned slide parts are not serialized (avoids duplicate part names).
sldIdLst = prs.slides._sldIdLst
for el in list(sldIdLst)[11:]:
    prs.part.drop_rel(el.get(qn("r:id")))
    sldIdLst.remove(el)

# --- Duplicate the clean bullet slide (idx 7) for the 4 topic slides ---
TITLE_ID, BODY_ID = 1293, 1294  # ids inherited from slide 7's template
topics = [
    # (title, bullets, image filename or None, image-layout)
    ("Dự đoán giá tiền mã hóa bằng học sâu", [
        "LSTM (Hochreiter & Schmidhuber, 1997) khắc phục vanishing gradient của RNN",
        "Ba cổng: input – forget – output điều phối cell state",
        "Fischer & Krauss (2018): LSTM vượt ARIMA & Random Forest",
        "Biến thể: Bi-LSTM, LSTM + Attention, TFT (Lim et al., 2021)",
    ], "lstm_architecture.png", "right"),
    ("Xử lý dữ liệu tài chính theo thời gian thực", [
        "Kafka + Spark Structured Streaming: bộ đôi phổ biến cho phân tích tài chính realtime",
        "Spark — unified engine cho cả batch lẫn streaming trên cùng API [Zaharia et al., 2016]",
        "Hỗ trợ cửa sổ trượt (sliding window) và watermark cho dữ liệu trễ",
        "Tích hợp MongoDB qua foreachBatch pattern, đảm bảo idempotency",
    ], None, None),
    ("Kiến trúc Lambda trong phân tích dữ liệu lớn", [
        "Nathan Marz (2015): tách biệt batch layer & speed layer, phục vụ qua serving layer",
        "Batch layer: độ chính xác cao, độ trễ cao (báo cáo, huấn luyện mô hình)",
        "Speed layer: độ trễ thấp, kết quả xấp xỉ (cảnh báo, dashboard realtime)",
        "Phù hợp hệ thống tài chính cần cả phân tích lịch sử & phản ứng nhanh",
    ], "architecture_overview.png", "bottom"),
    ("Chỉ số kỹ thuật & phân tích tài chính", [
        "Murphy (1999) hệ thống hóa phân tích kỹ thuật: SMA, RSI, Bollinger Bands, VWAP",
        "Các chỉ số được dùng rộng rãi và mang giá trị dự báo bổ sung",
        "Đề tài dùng các chỉ số này làm feature đầu vào cho LSTM thay vì raw price",
        "Kết hợp chỉ số kỹ thuật + học máy → tín hiệu dự báo mạnh hơn",
    ], None, None),
    ("Điểm mới của đề tài", [
        "Kết hợp Lambda Architecture với LSTM v3 (bổ sung volatility head)",
        "Direction-weighted Huber loss: ưu tiên dự đoán đúng hướng xu hướng",
        "Walk-forward validation (6 fold) + backtest 6 tháng — đánh giá thực tế",
        "Hệ thống end-to-end với Docker Compose, frontend React & model registry",
    ], None, None),
]

for title, bullets, img, mode in topics:
    ns = duplicate_slide(prs, 7)
    set_lines(get_shape(ns, TITLE_ID), [title])
    body = get_shape(ns, BODY_ID)
    set_lines(body, bullets)

    if img and mode == "right":
        # narrow the body to the left column, place image on the right
        body.left, body.width = Inches(0.79), Inches(4.6)
        path = os.path.join(FIG, img)
        w = Inches(4.1); h = Inches(4.1 / 1.85)
        ns.shapes.add_picture(path, Inches(5.55), Inches(2.35), w, h)
    elif img and mode == "bottom":
        # keep text in the upper area, place the wide diagram across the bottom
        body.top, body.height = Inches(1.30), Inches(1.55)
        path = os.path.join(FIG, img)
        w = Inches(7.6); h = Inches(7.6 / 3.01)
        ns.shapes.add_picture(path, Inches(1.20), Inches(3.05), w, h)

prs.save(OUT)
print("Saved:", OUT)
print("Total slides:", len(prs.slides))
