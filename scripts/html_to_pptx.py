"""
Convert crypto_presentation.html to a Google-Slides-compatible .pptx file
with editable text, matching the dark navy / gold design.
"""
import re
from pathlib import Path
from bs4 import BeautifulSoup, Tag
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1C, 0x26, 0x44)
GOLD   = RGBColor(0xC8, 0xA8, 0x70)
WARM   = RGBColor(0xE2, 0xDC, 0xD0)
MUTED  = RGBColor(0x8A, 0x96, 0xA8)
HINT   = RGBColor(0x4E, 0x5A, 0x6E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

# ── Slide dimensions: 16:9 widescreen ────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

MARGIN_L = Inches(1.1)
MARGIN_R = Inches(1.1)
MARGIN_T = Inches(0.5)
CONTENT_W = W - MARGIN_L - MARGIN_R


def rgb_bg(slide, color: RGBColor):
    """Fill slide background with solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height,
                text, font_name="Calibri", font_size=18, bold=False,
                color=WARM, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """Add a simple text box and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf


def add_paragraph(tf, text, font_name="Calibri", font_size=16,
                  bold=False, color=MUTED, indent=False, italic=False,
                  space_before=Pt(6)):
    """Append a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.space_before = space_before
    if indent:
        p.level = 1
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def gold_rule(slide, top):
    """Thin gold horizontal rule."""
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        MARGIN_L, top, Inches(0.4), Pt(1.5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()


def get_text(el) -> str:
    """Get clean text from a BS4 element."""
    if el is None:
        return ""
    return re.sub(r'\s+', ' ', el.get_text(separator=' ')).strip()


def bullets_from(ul_el) -> list[str]:
    """Extract list items from a <ul class="bullets"> or <ul class="checklist">."""
    if ul_el is None:
        return []
    return [get_text(li) for li in ul_el.find_all('li', recursive=False) if get_text(li)]


def table_rows(table_el) -> tuple[list[str], list[list[str]]]:
    """Return (headers, rows) from a sig-table element."""
    if table_el is None:
        return [], []
    headers = [get_text(th) for th in table_el.find_all('th')]
    rows = []
    for tr in table_el.find('tbody').find_all('tr'):
        rows.append([get_text(td) for td in tr.find_all('td')])
    return headers, rows


# ── Per-slide layout helpers ──────────────────────────────────────────────────

def layout_cover(prs, slide_data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    rgb_bg(slide, NAVY)

    top = Inches(1.8)

    # Kicker
    if slide_data.get('kicker'):
        add_textbox(slide, MARGIN_L, top, CONTENT_W, Inches(0.4),
                    slide_data['kicker'], font_name="Courier New",
                    font_size=11, color=GOLD)
        top += Inches(0.45)

    # Gold rule
    gold_rule(slide, top)
    top += Inches(0.15)

    # Title
    if slide_data.get('title'):
        add_textbox(slide, MARGIN_L, top, CONTENT_W, Inches(1.8),
                    slide_data['title'], font_name="Georgia",
                    font_size=40, bold=True, color=WARM)
        top += Inches(1.9)

    # Subtitle / italic
    if slide_data.get('subtitle'):
        add_textbox(slide, MARGIN_L, top, CONTENT_W, Inches(0.5),
                    slide_data['subtitle'], font_name="Georgia",
                    font_size=22, italic=True, color=GOLD)
        top += Inches(0.6)

    # Tags
    if slide_data.get('tags'):
        add_textbox(slide, MARGIN_L, top, CONTENT_W, Inches(0.4),
                    '  ·  '.join(slide_data['tags']),
                    font_name="Courier New", font_size=12, color=GOLD)
        top += Inches(0.5)

    # Stats row if present
    if slide_data.get('stats'):
        top += Inches(0.3)
        stat_w = CONTENT_W / len(slide_data['stats'])
        for i, (num, desc) in enumerate(slide_data['stats']):
            sx = MARGIN_L + stat_w * i
            add_textbox(slide, sx, top, stat_w - Inches(0.1), Inches(0.6),
                        num, font_name="Georgia", font_size=32,
                        bold=True, color=GOLD)
            add_textbox(slide, sx, top + Inches(0.55), stat_w - Inches(0.1),
                        Inches(0.5), desc, font_size=13, color=MUTED)

    # Bottom note
    if slide_data.get('note'):
        add_textbox(slide, MARGIN_L, H - Inches(1.2), CONTENT_W, Inches(0.5),
                    slide_data['note'], font_size=13, color=HINT)


def layout_content(prs, slide_data: dict):
    """Standard content slide: chrome bar + kicker + heading + body."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rgb_bg(slide, NAVY)

    # Chrome bar (top)
    chrome_h = Inches(0.55)
    chrome_box = slide.shapes.add_shape(1, 0, 0, W, chrome_h)
    chrome_box.fill.solid()
    chrome_box.fill.fore_color.rgb = RGBColor(0x23, 0x2F, 0x55)
    chrome_box.line.fill.background()

    section = slide_data.get('section', '')
    num = slide_data.get('num', '')
    if section:
        add_textbox(slide, MARGIN_L, Inches(0.1), CONTENT_W / 2, chrome_h,
                    section, font_name="Courier New", font_size=10, color=GOLD)
    if num:
        add_textbox(slide, W / 2, Inches(0.1), CONTENT_W / 2, chrome_h,
                    num, font_name="Courier New", font_size=10, color=MUTED,
                    align=PP_ALIGN.RIGHT)

    # Foot bar (bottom)
    foot_h = Inches(0.5)
    foot_box = slide.shapes.add_shape(1, 0, H - foot_h, W, foot_h)
    foot_box.fill.solid()
    foot_box.fill.fore_color.rgb = RGBColor(0x23, 0x2F, 0x55)
    foot_box.line.fill.background()
    add_textbox(slide, MARGIN_L, H - foot_h + Inches(0.08), CONTENT_W / 2,
                foot_h, "CS315 — Advanced Machine Learning",
                font_name="Courier New", font_size=9, color=HINT)
    add_textbox(slide, W / 2, H - foot_h + Inches(0.08), CONTENT_W / 2,
                foot_h, "Nhóm CN2 · 2026",
                font_name="Courier New", font_size=9, color=HINT,
                align=PP_ALIGN.RIGHT)

    top = chrome_h + Inches(0.3)

    # Kicker
    if slide_data.get('kicker'):
        add_textbox(slide, MARGIN_L, top, CONTENT_W, Inches(0.35),
                    slide_data['kicker'], font_name="Courier New",
                    font_size=10, color=GOLD)
        top += Inches(0.35)

    # Gold rule
    gold_rule(slide, top)
    top += Inches(0.18)

    # Heading
    if slide_data.get('heading'):
        h_box = add_textbox(slide, MARGIN_L, top, CONTENT_W, Inches(0.9),
                            slide_data['heading'], font_name="Georgia",
                            font_size=28, bold=True, color=WARM)
        top += Inches(0.95)

    # Body
    body_top = top
    body_h = H - foot_h - body_top - Inches(0.1)

    items = slide_data.get('body', [])
    if not items:
        return

    # If two-column, split evenly
    if slide_data.get('two_col'):
        col_w = (CONTENT_W - Inches(0.3)) / 2
        for ci, col in enumerate(slide_data['two_col']):
            cx = MARGIN_L + ci * (col_w + Inches(0.3))
            col_top = body_top
            if col.get('col_label'):
                add_textbox(slide, cx, col_top, col_w, Inches(0.3),
                            col['col_label'], font_name="Courier New",
                            font_size=9, color=GOLD)
                col_top += Inches(0.32)
            if col.get('bullets'):
                tf = add_textbox(slide, cx, col_top, col_w,
                                 body_h - Inches(0.32),
                                 '• ' + col['bullets'][0],
                                 font_size=14, color=MUTED)
                for b in col['bullets'][1:]:
                    add_paragraph(tf, '• ' + b, font_size=14, color=MUTED,
                                  space_before=Pt(8))
            elif col.get('table'):
                headers, rows = col['table']
                _add_table_text(slide, cx, col_top, col_w,
                                body_h - Inches(0.32), headers, rows)
            elif col.get('text'):
                add_textbox(slide, cx, col_top, col_w, body_h - Inches(0.32),
                            col['text'], font_size=14, color=MUTED)
        return

    # Table
    if slide_data.get('table'):
        headers, rows = slide_data['table']
        _add_table_text(slide, MARGIN_L, body_top, CONTENT_W, body_h,
                        headers, rows)
        if slide_data.get('note'):
            add_textbox(slide, MARGIN_L, H - foot_h - Inches(0.55),
                        CONTENT_W, Inches(0.4),
                        slide_data['note'], font_size=12, color=HINT)
        return

    # Stats row
    if slide_data.get('stats'):
        stat_w = CONTENT_W / len(slide_data['stats'])
        for i, (num_v, desc) in enumerate(slide_data['stats']):
            sx = MARGIN_L + stat_w * i
            add_textbox(slide, sx, body_top, stat_w - Inches(0.05),
                        Inches(0.65), num_v, font_name="Georgia",
                        font_size=30, bold=True, color=GOLD)
            add_textbox(slide, sx, body_top + Inches(0.65),
                        stat_w - Inches(0.05), Inches(0.5),
                        desc, font_size=13, color=MUTED)
        body_top += Inches(1.2)
        body_h -= Inches(1.2)

    # Bullets / checklist / plain
    if items:
        tf = add_textbox(slide, MARGIN_L, body_top, CONTENT_W, body_h,
                         '• ' + items[0], font_size=15, color=MUTED)
        for item in items[1:]:
            add_paragraph(tf, '• ' + item, font_size=15, color=MUTED,
                          space_before=Pt(10))

    if slide_data.get('note'):
        add_textbox(slide, MARGIN_L, H - foot_h - Inches(0.55),
                    CONTENT_W, Inches(0.4),
                    slide_data['note'], font_size=11, color=HINT)


def _add_table_text(slide, left, top, width, height, headers, rows):
    """Render a table as structured text boxes (PPTX native tables look great)."""
    if not rows and not headers:
        return
    n_cols = max(len(headers), max((len(r) for r in rows), default=0))
    n_rows = len(rows) + (1 if headers else 0)
    if n_cols == 0 or n_rows == 0:
        return

    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width,
                                 min(height, Inches(0.38 * n_rows + 0.1))).table
    tbl.columns[0].width = width // 3
    if n_cols > 1:
        rem = width - width // 3
        for ci in range(1, n_cols):
            tbl.columns[ci].width = rem // (n_cols - 1)

    row_idx = 0
    if headers:
        for ci, h in enumerate(headers[:n_cols]):
            cell = tbl.cell(0, ci)
            cell.text = h
            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = GOLD
            cell.text_frame.paragraphs[0].runs[0].font.bold = True
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].runs[0].font.name = "Courier New"
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x23, 0x2F, 0x55)
        row_idx = 1

    for r in rows:
        for ci, val in enumerate(r[:n_cols]):
            cell = tbl.cell(row_idx, ci)
            cell.text = val
            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = MUTED
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].runs[0].font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        row_idx += 1


# ── HTML parsing ──────────────────────────────────────────────────────────────

def parse_html(path: str) -> list[dict]:
    html = Path(path).read_text(encoding='utf-8')
    soup = BeautifulSoup(html, 'lxml')
    slides = []

    for section in soup.find_all('section', class_='slide'):
        d = {}

        # Section label + slide number from chrome
        chrome = section.find(class_='chrome')
        if chrome:
            labels = chrome.find_all(class_='label')
            if labels:
                d['section'] = get_text(labels[0])
            if len(labels) > 1:
                d['num'] = get_text(labels[1])

        # Cover slide?
        is_cover = 's-cover' in section.get('class', [])
        d['is_cover'] = is_cover

        # Kicker
        kicker = section.find(class_='kicker')
        if kicker:
            d['kicker'] = get_text(kicker)

        # Title / heading
        cover_title = section.find(class_='cover-title')
        if cover_title:
            d['title'] = get_text(cover_title)
        else:
            for cls in ['h2', 'h1', 'h3']:
                el = section.find(class_=cls)
                if el:
                    d['heading'] = get_text(el)
                    break

        # Cover subtitle
        cover_sub = section.find(class_='cover-sub')
        if cover_sub:
            d['subtitle'] = get_text(cover_sub)

        # Tags (cover)
        tags_el = section.find(class_='cover-tags')
        if tags_el:
            d['tags'] = [get_text(t) for t in tags_el.find_all(class_='tag')]

        # Lead / note
        lead = section.find(class_='lead')
        if lead:
            d['note'] = get_text(lead)

        # Two-column layout
        two_col_el = section.find(class_='two-col')
        if two_col_el:
            cols_raw = [c for c in two_col_el.children
                        if isinstance(c, Tag) and 'col-divider' not in c.get('class', [])]
            two_cols = []
            for col_el in cols_raw[:2]:
                col = {}
                kk = col_el.find(class_='kicker')
                if kk:
                    col['col_label'] = get_text(kk)
                ul = col_el.find('ul')
                if ul:
                    col['bullets'] = bullets_from(ul)
                tbl = col_el.find('table')
                if tbl:
                    col['table'] = table_rows(tbl)
                plain = col_el.find(class_=['body-text', 'lead', 'arch-desc', 'l-desc'])
                if plain and not col.get('bullets') and not col.get('table'):
                    col['text'] = get_text(plain)
                # arch-boxes
                arch_boxes = col_el.find_all(class_='arch-box')
                if arch_boxes and not col.get('bullets'):
                    combined = []
                    for ab in arch_boxes:
                        lbl = ab.find(class_='arch-label')
                        desc = ab.find(class_='arch-desc')
                        if lbl:
                            combined.append(get_text(lbl) + ':')
                        if desc:
                            combined.append(get_text(desc))
                    col['bullets'] = combined
                two_cols.append(col)
            if two_cols:
                d['two_col'] = two_cols

        # Table (top-level)
        if not d.get('two_col'):
            tbl = section.find('table')
            if tbl:
                d['table'] = table_rows(tbl)

        # Stats
        stat_row = section.find(class_='stat-row')
        if stat_row:
            stats = []
            for sc in stat_row.find_all(class_='stat-card'):
                num_el = sc.find(class_='stat-num')
                desc_el = sc.find(class_='stat-desc')
                stats.append((get_text(num_el), get_text(desc_el)))
            d['stats'] = stats

        # Big result
        big = section.find(class_='big-result')
        if big:
            num_el = big.find(class_='big-num')
            note_el = big.find(class_='big-note')
            lbl_el = big.find(class_='big-label')
            if num_el:
                d.setdefault('stats', []).insert(0,
                    (get_text(num_el), get_text(lbl_el) if lbl_el else ''))
            if note_el:
                d['note'] = get_text(note_el)

        # Bullets / checklist (top-level, outside two-col)
        if not d.get('two_col'):
            ul = section.find('ul')
            if ul:
                d['body'] = bullets_from(ul)

        # Evo-chain (LSTM evolution slide)
        evo = section.find(class_='evo-chain')
        if evo:
            steps = []
            for step in evo.find_all(class_='evo-step'):
                ver = get_text(step.find(class_='evo-ver'))
                name = get_text(step.find(class_='evo-name'))
                detail = get_text(step.find(class_='evo-detail'))
                issue = get_text(step.find(class_='evo-issue'))
                steps.append(f"{ver}: {name} — {detail}. {issue}")
            d['body'] = steps

        # Lambda layers
        lambda_wrap = section.find(class_='lambda-wrap')
        if lambda_wrap:
            items = []
            for layer in lambda_wrap.find_all(class_='lambda-layer'):
                lbl = get_text(layer.find(class_='l-label'))
                desc = get_text(layer.find(class_='l-desc'))
                items.append(f"[{lbl}] {desc}")
            d['body'] = items

        slides.append(d)

    return slides


# ── Main ──────────────────────────────────────────────────────────────────────

def convert(html_path: str, out_path: str):
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # Remove all default layouts' placeholder noise by using blank layout
    slide_data_list = parse_html(html_path)

    for sd in slide_data_list:
        if sd.get('is_cover'):
            layout_cover(prs, sd)
        else:
            layout_content(prs, sd)

    prs.save(out_path)
    print(f"Saved: {out_path}  ({len(slide_data_list)} slides)")


if __name__ == '__main__':
    import sys
    html_in = sys.argv[1] if len(sys.argv) > 1 else \
        "/Users/ducqhle/Documents/workspace/CS315.F21.CN2.AdvanceML/crypto_presentation.html"
    pptx_out = sys.argv[2] if len(sys.argv) > 2 else \
        html_in.replace('.html', '.pptx')
    convert(html_in, pptx_out)
